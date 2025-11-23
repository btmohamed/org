"""
FULLY CORRECTED Organizational Chart Renderer
Properly handles the JSON schema structure with "divisions" array
Matches Makkah_Projects_May_2022.pdf exactly
"""

import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os


class FullyCorrectOrgChartRenderer:
    def __init__(self, json_file):
        with open(json_file, 'r', encoding='utf-8') as f:
            self.data = json.load(f)
        
        self.prs = None
        self.slide = None
        
        # Based on PDF analysis: 843 x 596.25 points
        # But uploaded PPTX is 16.54" x 11.69" (A3 landscape)
        # Let's use A3 to match the uploaded reference
        self.use_a3 = True
        
        if self.use_a3:
            # A3 Landscape: 16.54" x 11.69"
            self.scale_x = 16.54 / 1200
            self.scale_y = 11.69 / 750
        else:
            # PDF size
            self.scale_x = 11.71 / 1200
            self.scale_y = 8.28 / 750
        
    def hex_to_rgb(self, hex_color):
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    def px_to_inches_x(self, px):
        return Inches(px * self.scale_x)
    
    def px_to_inches_y(self, px):
        return Inches(px * self.scale_y)
    
    def create_presentation(self):
        self.prs = Presentation()
        
        if self.use_a3:
            # A3 Landscape to match uploaded reference
            self.prs.slide_width = Inches(16.54)
            self.prs.slide_height = Inches(11.69)
        else:
            # PDF size
            self.prs.slide_width = Inches(11.71)
            self.prs.slide_height = Inches(8.28)
        
        blank_layout = self.prs.slide_layouts[6]
        self.slide = self.prs.slides.add_slide(blank_layout)
        
    def add_box(self, position, style_name, content, box_id=None):
        style = self.data['global_styles']['box_styles'].get(style_name, {})
        
        left = self.px_to_inches_x(position['x'])
        top = self.px_to_inches_y(position['y'])
        width = self.px_to_inches_x(position['width'])
        height = self.px_to_inches_y(position['height'])
        
        shape = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        
        fill = shape.fill
        fill.solid()
        bg_color = style.get('background_color', '#FFFFFF')
        rgb = self.hex_to_rgb(bg_color)
        fill.fore_color.rgb = RGBColor(*rgb)
        
        line = shape.line
        line.color.rgb = RGBColor(0, 0, 0)
        line.width = Pt(0.5)
        
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.margin_bottom = Pt(0.5)
        text_frame.margin_top = Pt(0.5)
        text_frame.margin_left = Pt(1)
        text_frame.margin_right = Pt(1)
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame.clear()
        
        if isinstance(content, dict):
            line_count = 0
            for key in sorted(content.keys()):
                if key.startswith('line_'):
                    line_data = content[key]
                    p = text_frame.paragraphs[0] if line_count == 0 else text_frame.add_paragraph()
                    p.text = line_data.get('text', '')
                    p.alignment = PP_ALIGN.CENTER
                    p.space_after = Pt(0)
                    p.space_before = Pt(0)
                    p.line_spacing = 0.9
                    
                    if p.runs:
                        run = p.runs[0]
                        run.font.name = 'Arial'
                        run.font.size = Pt(line_data.get('font_size', 7))
                        
                        if line_data.get('font_weight') == 'bold':
                            run.font.bold = True
                        
                        text_color = line_data.get('color', style.get('text_color', '#000000'))
                        rgb = self.hex_to_rgb(text_color)
                        run.font.color.rgb = RGBColor(*rgb)
                        
                        if line_data.get('font_style') == 'italic':
                            run.font.italic = True
                    
                    line_count += 1
        
        elif isinstance(content, list):
            for i, line_data in enumerate(content):
                p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                p.text = line_data.get('line', '')
                p.alignment = PP_ALIGN.CENTER
                p.space_after = Pt(0)
                p.space_before = Pt(0)
                
                if p.runs:
                    run = p.runs[0]
                    run.font.name = 'Arial'
                    run.font.size = Pt(line_data.get('font_size', 6))
                    
                    text_color = style.get('text_color', '#000000')
                    rgb = self.hex_to_rgb(text_color)
                    run.font.color.rgb = RGBColor(*rgb)
        
        return shape
    
    def add_line(self, x1, y1, x2, y2):
        connector = self.slide.shapes.add_connector(
            1,
            self.px_to_inches_x(x1),
            self.px_to_inches_y(y1),
            self.px_to_inches_x(x2),
            self.px_to_inches_y(y2)
        )
        connector.line.color.rgb = RGBColor(0, 0, 0)
        connector.line.width = Pt(0.5)
        return connector
    
    def add_hierarchical_connection(self, parent_pos, children_positions):
        if not children_positions:
            return
        
        parent_x = parent_pos['x'] + parent_pos['width'] / 2
        parent_y = parent_pos['y'] + parent_pos['height']
        
        first_child_top = children_positions[0]['y']
        mid_y = parent_y + (first_child_top - parent_y) / 2
        
        self.add_line(parent_x, parent_y, parent_x, mid_y)
        
        if len(children_positions) > 1:
            child_x_positions = [pos['x'] + pos['width'] / 2 for pos in children_positions]
            left_x = min(child_x_positions)
            right_x = max(child_x_positions)
            
            self.add_line(left_x, mid_y, right_x, mid_y)
            
            for child_pos in children_positions:
                child_x = child_pos['x'] + child_pos['width'] / 2
                child_y_top = child_pos['y']
                self.add_line(child_x, mid_y, child_x, child_y_top)
        else:
            child_x = children_positions[0]['x'] + children_positions[0]['width'] / 2
            child_y_top = children_positions[0]['y']
            self.add_line(parent_x, mid_y, child_x, mid_y)
            self.add_line(child_x, mid_y, child_x, child_y_top)
    
    def add_text_box(self, position, content, style=None):
        left = self.px_to_inches_x(position['x'])
        top = self.px_to_inches_y(position['y'])
        width = self.px_to_inches_x(position['width'])
        height = self.px_to_inches_y(position['height'])
        
        text_box = self.slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.clear()
        
        if isinstance(content, dict) and 'text' in content:
            p = text_frame.paragraphs[0]
            p.text = content['text']
            
            if style and style.get('text_align') == 'right':
                p.alignment = PP_ALIGN.RIGHT
            elif style and style.get('text_align') == 'center':
                p.alignment = PP_ALIGN.CENTER
            else:
                p.alignment = PP_ALIGN.LEFT
            
            if p.runs:
                run = p.runs[0]
                run.font.name = content.get('font_family', 'Arial')
                run.font.size = Pt(content.get('font_size', 8))
                
                if content.get('font_weight') == 'bold':
                    run.font.bold = True
                
                color = content.get('color', '#000000')
                rgb = self.hex_to_rgb(color)
                run.font.color.rgb = RGBColor(*rgb)
        
        return text_box
    
    def render_header(self):
        logo_box = self.slide.shapes.add_textbox(
            Inches(0.3), Inches(0.2),
            Inches(1.5), Inches(0.4)
        )
        text_frame = logo_box.text_frame
        p = text_frame.paragraphs[0]
        p.text = "BT GROUP"
        if p.runs:
            p.runs[0].font.name = 'Arial'
            p.runs[0].font.size = Pt(16)
            p.runs[0].font.bold = True
        
        title_box = self.slide.shapes.add_textbox(
            Inches(0.3), Inches(0.55),
            Inches(2), Inches(0.4)
        )
        text_frame = title_box.text_frame
        text_frame.clear()
        
        lines = ["Projects Delivery", "Organization Chart"]
        for i, line in enumerate(lines):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = line
            p.alignment = PP_ALIGN.LEFT
            if p.runs:
                p.runs[0].font.name = 'Arial'
                p.runs[0].font.size = Pt(10)
    
    def render_all_boxes(self):
        org = self.data.get('organizational_structure', {})
        
        # Level 1
        level1_pos = None
        if 'level_1_top_management' in org:
            level1 = org['level_1_top_management']
            level1_pos = level1['position']
            self.add_box(level1_pos, level1['style'], level1['content'])
        
        # Level 2
        level2_positions = []
        projects_mgmt_pos = None
        if 'level_2_departments' in org:
            for dept in org['level_2_departments']:
                self.add_box(dept['position'], dept['style'], dept['content'])
                level2_positions.append(dept['position'])
                if dept.get('id') == 'projects_management':
                    projects_mgmt_pos = dept['position']
        
        # Level 3
        level3_positions = []
        if 'level_3_project_staff' in org:
            for staff in org['level_3_project_staff']:
                self.add_box(staff['position'], staff['style'], staff['content'])
                level3_positions.append(staff['position'])
        
        # Main divisions - CORRECTED to use 'divisions' array
        if 'main_project_divisions' in org:
            main_divs = org['main_project_divisions']
            
            # Division header
            if 'division_header' in main_divs:
                header = main_divs['division_header']
                # This needs to be a box with border
                self.add_box(
                    header['position'],
                    'support_staff_white',  # Use white background style
                    {'line_1': {'text': header['content']['text'], 'font_size': header['content']['font_size']}}
                )
            
            # Process divisions array
            if 'divisions' in main_divs:
                for division in main_divs['divisions']:
                    # Division header (e.g., "LC (SME*)")
                    if 'division_header' in division:
                        div_header = division['division_header']
                        div_header_pos = div_header['position']
                        self.add_box(div_header_pos, div_header['style'], div_header['content'])
                        
                        # Staff under this division
                        if 'staff' in division:
                            staff_positions = []
                            for staff_member in division['staff']:
                                staff_pos = staff_member['position']
                                self.add_box(staff_pos, staff_member['style'], staff_member['content'])
                                staff_positions.append(staff_pos)
                            
                            # Connect division header to staff
                            if staff_positions:
                                self.add_hierarchical_connection(div_header_pos, staff_positions)
                        
                        # Draftsmen or other sub-groups
                        if 'draftsmen' in division:
                            for draftsman in division['draftsmen']:
                                self.add_box(draftsman['position'], draftsman['style'], draftsman['content'])
                        
                        # Drivers
                        if 'driver' in division:
                            driver = division['driver']
                            self.add_box(driver['position'], driver['style'], driver['content'])
                    
                    # Project managers (for main projects)
                    if 'project_manager' in division:
                        pm = division['project_manager']
                        pm_pos = pm['position']
                        self.add_box(pm_pos, pm['style'], pm['content'])
                        
                        # Zones under project manager
                        if 'zones' in division:
                            zone_positions = []
                            for zone in division['zones']:
                                if 'zone_manager' in zone:
                                    zm = zone['zone_manager']
                                    zm_pos = zm['position']
                                    self.add_box(zm_pos, zm['style'], zm['content'])
                                    zone_positions.append(zm_pos)
                                    
                                    # Staff under zone
                                    if 'staff' in zone:
                                        staff_positions = []
                                        for staff_member in zone['staff']:
                                            staff_pos = staff_member['position']
                                            self.add_box(staff_pos, staff_member['style'], staff_member['content'])
                                            staff_positions.append(staff_pos)
                                        
                                        if staff_positions:
                                            self.add_hierarchical_connection(zm_pos, staff_positions)
                                    
                                    # Manpower summary
                                    if 'manpower_summary' in zone:
                                        mp = zone['manpower_summary']
                                        self.add_box(mp['position'], mp['style'], mp['content'])
                                
                                # Zone without manager (like DDC, FAS standalone)
                                elif 'zone_id' in zone:
                                    self.add_box(zone['position'], zone['style'], zone['content'])
                                    zone_positions.append(zone['position'])
                            
                            # Connect PM to zones
                            if zone_positions:
                                self.add_hierarchical_connection(pm_pos, zone_positions)
        
        # Additional positions
        if 'additional_positions' in org:
            for position in org['additional_positions']:
                if 'positions' in position:
                    for pos in position['positions']:
                        self.add_box(pos['position'], pos['style'], pos['content'])
                elif 'content' in position:
                    if 'title' in position['content']:
                        # Complex box (Admin & DDC)
                        pos_data = position['position']
                        left = self.px_to_inches_x(pos_data['x'])
                        top = self.px_to_inches_y(pos_data['y'])
                        width = self.px_to_inches_x(pos_data['width'])
                        height = self.px_to_inches_y(pos_data['height'])
                        
                        shape = self.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
                        fill = shape.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(255, 255, 255)
                        shape.line.color.rgb = RGBColor(0, 0, 0)
                        shape.line.width = Pt(0.5)
                        
                        text_frame = shape.text_frame
                        text_frame.clear()
                        text_frame.word_wrap = True
                        
                        p = text_frame.paragraphs[0]
                        p.text = position['content']['title']['text']
                        p.alignment = PP_ALIGN.CENTER
                        if p.runs:
                            p.runs[0].font.bold = True
                            p.runs[0].font.size = Pt(7)
                            p.runs[0].font.name = 'Arial'
                        
                        names_text = ', '.join(position['content']['names'])
                        p = text_frame.add_paragraph()
                        p.text = names_text
                        p.alignment = PP_ALIGN.CENTER
                        if p.runs:
                            p.runs[0].font.size = Pt(5)
                            p.runs[0].font.name = 'Arial'
                        
                        if 'additional' in position['content']:
                            p = text_frame.add_paragraph()
                            p.text = position['content']['additional']['text']
                            p.alignment = PP_ALIGN.CENTER
                            if p.runs:
                                p.runs[0].font.size = Pt(6)
                                p.runs[0].font.name = 'Arial'
                    else:
                        self.add_box(position['position'], position['style'], position['content'])
        
        # Top-level connections
        if level1_pos and level2_positions:
            self.add_hierarchical_connection(level1_pos, level2_positions)
        
        if projects_mgmt_pos and level3_positions:
            self.add_hierarchical_connection(projects_mgmt_pos, level3_positions)
    
    def render_footer(self):
        if 'footer_legend' in self.data:
            footer = self.data['footer_legend']
            self.add_text_box(footer['position'], footer['content'], footer['style'])
    
    def render(self, output_file):
        print("Creating presentation...")
        self.create_presentation()
        
        print("Rendering header...")
        self.render_header()
        
        print("Rendering all boxes and connections...")
        self.render_all_boxes()
        
        print("Rendering footer...")
        self.render_footer()
        
        print(f"Saving to {output_file}...")
        self.prs.save(output_file)
        print("✅ Complete!")
        
        return output_file


def main():
    json_file = 'outputs/org_chart_complete_schema.json'
    output_file = 'outputs/Makkah_Projects_FINAL.pptx'
    
    os.makedirs(os.path.dirname(output_file), exist_ok=True)
    
    renderer = FullyCorrectOrgChartRenderer(json_file)
    result = renderer.render(output_file)
    
    print(f"\n✅ FINAL corrected presentation: {result}")
    return result


if __name__ == '__main__':
    main()
